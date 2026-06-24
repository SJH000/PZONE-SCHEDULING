from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("outputs/2026_06_14")
RULE_BASE = BASE / "rule_baseline"
RULE_REPLAY = BASE / "rule_replay"
RULE_SIM = BASE / "rule_simulation"
OUT = BASE / "PZONE_rule_baseline_replay_7slides.pptx"
SCRIPT_OUT = BASE / "PZONE_rule_baseline_replay_ppt_script_1min.md"

FONT = "Malgun Gothic"
NAVY = RGBColor(24, 39, 61)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(185, 28, 28)
GRAY = RGBColor(75, 85, 99)
LIGHT_GRAY = RGBColor(243, 244, 246)


def set_text_style(shape, font_size=14, color=NAVY, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, font_size=14, color=NAVY, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    set_text_style(shape, font_size, color, bold)
    return shape


def add_title(slide, title, subtitle=None, section=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_textbox(slide, 0.45, 0.35, 10.5, 0.45, title, 23, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.48, 0.82, 11.0, 0.32, subtitle, 10.5, GRAY)
    if section:
        add_textbox(slide, 11.65, 0.42, 1.2, 0.3, section, 9.5, GRAY, False, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, font_size=12, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_callout(slide, x, y, w, h, title, body, color=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(239, 246, 255)
    rect.line.color.rgb = color
    rect.line.width = Pt(1)
    add_textbox(slide, x + 0.16, y + 0.1, w - 0.32, 0.24, title, 10, color, True)
    add_textbox(slide, x + 0.16, y + 0.38, w - 0.32, h - 0.48, body, 11, NAVY)


def add_table(slide, x, y, w, h, headers, rows, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = str(value)
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = NAVY
    for cell in table.rows[0].cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
    return table_shape


def add_image_fit(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio >= box_ratio:
        final_w = w
        final_h = w / ratio
    else:
        final_h = h
        final_w = h * ratio
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))


def add_footer(slide, page):
    add_textbox(slide, 0.45, 7.12, 5.3, 0.2, "P-ZONE Scheduling | Rule baseline & offline evaluation", 8.5, GRAY)
    add_textbox(slide, 12.3, 7.12, 0.55, 0.2, f"{page}/7", 8.5, GRAY, False, PP_ALIGN.RIGHT)


def add_notes(slide, note):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = note
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(10)


def fmt(value, digits=1):
    try:
        return f"{float(value):,.{digits}f}"
    except Exception:
        return str(value)


def pct(value):
    try:
        return f"{float(value) * 100:.1f}%"
    except Exception:
        return str(value)


def load_data():
    counts = pd.read_csv(RULE_BASE / "data" / "rule_trigger_counts.csv")
    replay = pd.read_csv(RULE_REPLAY / "data" / "replay_effect_summary.csv")
    quality = pd.read_csv(RULE_REPLAY / "data" / "replay_action_quality.csv")
    sim = pd.read_csv(RULE_SIM / "data" / "simulation_effect_summary.csv")
    state = pd.read_csv(RULE_BASE / "data" / "state_5min.csv")
    actions = pd.read_csv(RULE_BASE / "data" / "actions.csv")
    return counts, replay, quality, sim, state, actions


def get_replay(replay, action, metric, col):
    row = replay[(replay["action"] == action) & (replay["metric"] == metric)]
    if row.empty:
        return None
    return row.iloc[0][col]


def build():
    counts, replay, quality, sim, state, actions = load_data()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes = {}

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "6~7번 진행 목적", "병목 진단 이후 rule 후보를 만들고, 기존 로그로 적용 가능성을 오프라인 평가", "01")
    add_callout(slide, 0.6, 1.2, 5.9, 1.05, "전제", "현장 추가정보가 없으므로 실제 설비 제어가 아니라 로그 기반 rule 설계와 proxy 검증으로 진행했다.", ORANGE)
    add_table(
        slide,
        0.65,
        2.75,
        5.8,
        2.35,
        ["단계", "목적", "산출물"],
        [
            ["6번", "Rule baseline 설계", "state, actions, trigger 검증"],
            ["7번", "Replay 평가", "trigger 전후 60분 비교"],
            ["7번", "Counterfactual", "10/20/30% 완화 시나리오"],
        ],
        9,
    )
    add_bullets(
        slide,
        7.0,
        1.25,
        5.4,
        2.0,
        [
            "입력: 1~5번 분석에서 만든 정상 route, WIP, 대기, rail state",
            "출력: action 후보와 trigger 시점, 전후 변화, 완화 민감도",
            "해석: 실제 개선 확정이 아니라 다음 설계 단계의 근거 자료",
        ],
        12.2,
    )
    add_callout(slide, 7.0, 4.15, 5.35, 1.0, "핵심 메시지", "A3 우선 반출 rule은 가장 먼저 검증할 후보이고, A4/Rail rule은 현장 매핑 전까지 보조 후보로 해석한다.", BLUE)
    notes[1] = "이 장은 6번과 7번의 목적을 설명하는 장입니다. 1~5번에서 병목 위치를 찾았고, 6번부터는 개선 가능성을 보기 위한 rule baseline을 만들었습니다. 다만 현장 추가정보가 없기 때문에 실제 설비 제어가 아니라 기존 로그 기반의 오프라인 평가입니다. 6번은 어떤 상황에서 어떤 action을 낼지 정했고, 7번은 그 action이 발생한 시점 전후로 병목 지표가 어떻게 변했는지와 완화 가정 시 어떤 지표가 민감한지 확인했습니다."
    add_notes(slide, notes[1])
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Rule-based Scheduling Baseline 설계", "A3, A4, Rail/AMR 병목 후보를 action rule로 변환", "02")
    rule_rows = [
        ["PRIORITIZE_DISCHARGE", "BUFFER WIP >= 70 또는 A3 후단 대기 증가", "A3/BUFFER/C2", "High"],
        ["DISPATCH_AMR", "A3 완료품 + A3->AMR 대기 증가", "AMR_LD250", "Medium"],
        ["CONTROL_WIP", "A4 long window 발생", "A4 WIP", "Medium"],
        ["HOLD_ENTRY", "A4 long + rail congestion 동시 발생", "A4 entry", "Medium"],
        ["SCHEDULE_RAIL", "rail event p90 이상 또는 active rail 4", "TR01", "Low"],
    ]
    add_table(slide, 0.55, 1.2, 12.3, 2.55, ["Action", "Trigger 조건", "대상", "신뢰도"], rule_rows, 8.2)
    add_image_fit(slide, RULE_BASE / "figures" / "rule_trigger_counts.png", 0.7, 4.0, 5.7, 2.35)
    add_bullets(
        slide,
        6.9,
        4.05,
        5.7,
        1.75,
        [
            "High: A3 후단 대기/WIP처럼 로그 근거가 강한 action",
            "Medium: AMR 가능 여부나 A4 진입 제어처럼 추가 제약이 필요한 action",
            "Low: TRN_DEV_ID 물리 위치가 없어 구간 확정이 어려운 rail action",
        ],
        11,
    )
    notes[2] = "두 번째 장은 rule baseline의 설계 내용입니다. 핵심 action은 다섯 가지입니다. A3 후단 대기나 버퍼 WIP가 높으면 PRIORITIZE_DISCHARGE를 내고, A3 완료품과 AMR 대기가 같이 보이면 DISPATCH_AMR을 냅니다. A4 long window가 있으면 CONTROL_WIP, A4 long과 rail 혼잡이 같이 있으면 HOLD_ENTRY, rail 이벤트가 과도하면 SCHEDULE_RAIL을 냅니다. 단, AMR 가능 여부와 rail 물리 위치 정보가 없기 때문에 action마다 신뢰도를 따로 두었습니다."
    add_notes(slide, notes[2])
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "State / Action Log 생성 결과", "5분 bucket 단위 state를 만들고 rule을 적용해 action log 생성", "03")
    count_rows = []
    for _, r in counts.iterrows():
        count_rows.append([r["action"], int(r["trigger_count"]), int(r["bucket_count"]), r.get("confidence_modes", "")])
    add_table(slide, 0.55, 1.18, 5.55, 2.5, ["Action", "Trigger", "Bucket", "Confidence"], count_rows, 8.3)
    add_image_fit(slide, RULE_BASE / "figures" / "rule_trigger_timeline.png", 6.45, 1.15, 6.25, 3.05)
    add_callout(
        slide,
        0.65,
        4.15,
        5.35,
        1.35,
        "생성 규모",
        f"state bucket {len(state):,}개, action row {len(actions):,}개를 생성했다. NO_ACTION bucket은 {int(counts[counts['action']=='NO_ACTION']['bucket_count'].iloc[0]):,}개다.",
        BLUE,
    )
    add_bullets(
        slide,
        6.75,
        4.55,
        5.65,
        1.1,
        [
            "PRIORITIZE_DISCHARGE가 가장 많이 발생: A3 후단 pressure를 넓게 감지",
            "SCHEDULE_RAIL은 rail event p90 기준으로 제한해 과도한 trigger를 방지",
            "DISPATCH_AMR은 AMR availability 부재로 보수적으로 적게 발생",
        ],
        10.8,
    )
    notes[3] = "세 번째 장은 state와 action log 생성 결과입니다. 기존 로그를 5분 단위로 재구성해서 2만 4천여 개 state bucket을 만들었고, 그 위에 rule을 적용해 1만 2천여 개 action row를 생성했습니다. 가장 많이 발생한 것은 PRIORITIZE_DISCHARGE입니다. 이는 A3 후단 압박과 BUFFER WIP를 넓게 잡기 때문입니다. SCHEDULE_RAIL은 rail event p90 이상일 때만 발생하도록 조정했고, DISPATCH_AMR은 AMR availability가 없어서 보수적으로 적게 발생했습니다."
    add_notes(slide, notes[3])
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Replay 평가: A3 후단 반출 Rule", "Action trigger 전후 60분의 A3 대기와 BUFFER WIP 변화 관찰", "04")
    add_image_fit(slide, RULE_REPLAY / "figures" / "a3_before_after_wait.png", 0.55, 1.1, 6.0, 3.1)
    add_image_fit(slide, RULE_REPLAY / "figures" / "buffer_wip_after_rule.png", 6.75, 1.1, 5.75, 3.1)
    rows = [
        ["PRIORITIZE_DISCHARGE", "A3->AMR", fmt(get_replay(replay, "PRIORITIZE_DISCHARGE", "a3_to_amr_wait_p90_recent", "delta_mean")), "소폭 감소"],
        ["PRIORITIZE_DISCHARGE", "A3->C2", fmt(get_replay(replay, "PRIORITIZE_DISCHARGE", "a3_to_c2_wait_p90_recent", "delta_mean")), "증가"],
        ["PRIORITIZE_DISCHARGE", "BUFFER WIP", fmt(get_replay(replay, "PRIORITIZE_DISCHARGE", "buffer_wip", "delta_mean")), "거의 변화 없음"],
        ["DISPATCH_AMR", "A3->AMR", fmt(get_replay(replay, "DISPATCH_AMR", "a3_to_amr_wait_p90_recent", "delta_mean")), "감소 경향"],
    ]
    add_table(slide, 0.75, 4.55, 11.7, 1.25, ["Action", "Metric", "After-Before", "관찰"], rows, 8.8)
    add_callout(slide, 0.75, 6.05, 11.7, 0.45, "해석", "A3 rule은 A3->AMR 대기 완화 가능성을 보이지만, C2 저장 대기와 BUFFER WIP까지 동시에 개선됐다고 보기는 어렵다.", ORANGE)
    notes[4] = "네 번째 장은 A3 후단 반출 rule의 replay 평가입니다. 그래프는 action 발생 전 60분과 후 60분을 비교한 것입니다. PRIORITIZE_DISCHARGE는 A3에서 AMR로 가는 대기는 소폭 낮아지는 경향이 있지만, A3에서 C2로 이어지는 저장 대기와 BUFFER WIP는 뚜렷하게 개선됐다고 보기 어렵습니다. DISPATCH_AMR은 A3->AMR 대기는 줄어드는 방향이지만 trigger 수가 14건으로 적고, AMR 실제 가용 여부가 없기 때문에 확정할 수는 없습니다. 따라서 A3 rule은 우선 검증 후보지만 후단 C2와 buffer 문제를 함께 봐야 합니다."
    add_notes(slide, notes[4])
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Replay 평가: A4 / Rail Rule", "A4 long window와 rail congestion action의 전후 변화 관찰", "05")
    add_image_fit(slide, RULE_REPLAY / "figures" / "a4_before_after_long.png", 0.55, 1.1, 6.0, 3.15)
    add_image_fit(slide, RULE_REPLAY / "figures" / "rail_before_after_congestion.png", 6.75, 1.1, 5.75, 3.15)
    rows = [
        ["CONTROL_WIP", "A4 long count", fmt(get_replay(replay, "CONTROL_WIP", "a4_long_count", "delta_mean"), 3), "소폭 감소"],
        ["CONTROL_WIP", "A4 processing max", fmt(get_replay(replay, "CONTROL_WIP", "a4_processing_max_sec", "delta_mean")), "소폭 감소"],
        ["HOLD_ENTRY", "A4 metrics", "증가", "자연 개선 신호 약함"],
        ["SCHEDULE_RAIL", "rail event count", fmt(get_replay(replay, "SCHEDULE_RAIL", "rail_event_count", "delta_mean")), "증가"],
    ]
    add_table(slide, 0.75, 4.55, 11.7, 1.25, ["Action", "Metric", "After-Before", "관찰"], rows, 8.8)
    add_callout(slide, 0.75, 6.05, 11.7, 0.45, "해석", "A4 CONTROL_WIP는 소폭 개선 신호가 있으나, HOLD_ENTRY와 Rail scheduling은 현장 물리 매핑 없이는 확정성이 낮다.", ORANGE)
    notes[5] = "다섯 번째 장은 A4와 Rail rule의 replay 결과입니다. CONTROL_WIP는 A4 long count와 A4 processing max가 소폭 낮아지는 관찰 결과가 있습니다. 하지만 HOLD_ENTRY는 trigger 이후 오히려 A4 지표가 증가하는 경향이 있어, 현재 rule만으로는 충분하지 않을 수 있습니다. SCHEDULE_RAIL도 trigger 이후 rail event count가 증가하는 관찰 결과가 있어 자연 개선 신호가 강하지 않습니다. 특히 rail은 TRN_DEV_ID의 물리 위치가 없기 때문에 실제 어느 구간을 조정해야 하는지 확정할 수 없습니다."
    add_notes(slide, notes[5])
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Counterfactual 시뮬레이션 결과", "Action이 10/20/30% 완화 효과를 낸다고 가정한 민감도 분석", "06")
    add_image_fit(slide, RULE_SIM / "figures" / "scenario_wait_reduction.png", 0.45, 1.0, 6.05, 2.75)
    add_image_fit(slide, RULE_SIM / "figures" / "scenario_wip_reduction.png", 6.85, 1.0, 5.65, 2.75)
    add_image_fit(slide, RULE_SIM / "figures" / "scenario_action_comparison.png", 1.5, 4.0, 10.3, 2.15)
    add_callout(slide, 0.75, 6.35, 11.7, 0.38, "해석", "A3 후단 대기와 BUFFER WIP는 완화율에 민감하지만, 10/20/30%는 가정값이므로 실제 개선율로 말하면 안 된다.", BLUE)
    notes[6] = "여섯 번째 장은 counterfactual 시뮬레이션입니다. 여기서는 실제 설비가 제어된 것이 아니기 때문에, action이 효과적이었다고 가정하고 after 지표를 10%, 20%, 30% 줄여봤습니다. 이 분석은 실제 개선율을 말하는 것이 아니라, 어떤 지표가 완화율에 민감한지를 보는 것입니다. 결과적으로 A3 후단 대기와 BUFFER WIP는 완화 가정에 민감하게 반응합니다. 즉 A3 우선 반출이나 AMR/C2 연계 개선이 실제로 가능하다면 효과가 날 여지가 큽니다. 다만 이 수치는 가정 기반이므로 실제 제어 효과로 표현하면 안 됩니다."
    add_notes(slide, notes[6])
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "현재 판단과 다음 단계", "Rule baseline은 개선 후보 선별용으로 유효, 다음은 rule 정교화와 LLM agent 설계", "07")
    rows = [
        ["A3 우선 반출", "가장 먼저 검증할 후보", "A3->AMR 일부 감소, WIP/C2는 추가 분리 필요"],
        ["AMR dispatch", "보조 후보", "trigger 적고 availability 부재"],
        ["A4 WIP control", "보조 개선 후보", "소폭 개선 관찰"],
        ["HOLD_ENTRY", "rule 재조정 필요", "replay상 자연 개선 약함"],
        ["Rail scheduling", "현장정보 필요", "TRN_DEV_ID 물리 위치 부재"],
    ]
    add_table(slide, 0.55, 1.15, 12.3, 2.25, ["영역", "판단", "근거"], rows, 8.5)
    add_bullets(
        slide,
        0.8,
        3.85,
        5.75,
        1.85,
        [
            "다음 rule 개선: PRIORITIZE_DISCHARGE 조건을 WIP 단독보다 후단 대기와 결합",
            "A4: HOLD_ENTRY 조건을 rail congestion보다 A4 주변 WIP/long window 중심으로 재조정",
            "Rail: 물리 위치 없이는 dispatch 후보만 유지",
        ],
        11,
    )
    add_bullets(
        slide,
        7.0,
        3.85,
        5.6,
        1.85,
        [
            "LLM agent 입력: 5분 state + action candidates + constraints",
            "출력: action, target, reason, confidence, limitation",
            "검증: rule baseline과 같은 replay/simulation framework로 비교",
        ],
        11,
    )
    add_callout(slide, 0.75, 6.05, 11.7, 0.55, "결론", "6~7번 결과는 A3 중심 rule을 우선 개선 대상으로 남기고, A4/Rail은 현장 제약 정보가 보강될 때 정교화하는 방향을 제시한다.", TEAL)
    notes[7] = "마지막 장은 현재 판단과 다음 단계입니다. 6번과 7번 결과를 보면, A3 우선 반출 rule은 가장 먼저 검증할 후보입니다. 다만 BUFFER WIP만으로 너무 넓게 trigger되는 문제가 있으므로 후단 대기 조건과 결합해서 정교화해야 합니다. A4 CONTROL_WIP는 보조 개선 후보로 유지할 수 있지만, HOLD_ENTRY는 현재 rule을 재조정할 필요가 있습니다. Rail scheduling은 물리 위치 정보가 없어서 구체 제어로 가기 어렵습니다. 다음 단계는 이 rule baseline을 기준으로 LLM decision agent의 state, action schema, constraint verifier를 설계하고 같은 replay framework로 비교 평가하는 것입니다."
    add_notes(slide, notes[7])
    add_footer(slide, 7)

    prs.save(OUT)
    write_script(notes)
    print(f"saved {OUT}")
    print(f"saved {SCRIPT_OUT}")


def write_script(notes):
    titles = [
        "6~7번 진행 목적",
        "Rule-based Scheduling Baseline 설계",
        "State / Action Log 생성 결과",
        "Replay 평가: A3 후단 반출 Rule",
        "Replay 평가: A4 / Rail Rule",
        "Counterfactual 시뮬레이션 결과",
        "현재 판단과 다음 단계",
    ]
    lines = ["# P-ZONE 6~7번 진행 공유 PPT 1분 대본", ""]
    for i, title in enumerate(titles, start=1):
        lines.append(f"## {i}. {title}")
        lines.append("")
        lines.append(notes[i])
        lines.append("")
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    build()
