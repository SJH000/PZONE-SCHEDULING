from __future__ import annotations

import math
import sqlite3
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(".")
SRC_DATA = ROOT / "outputs" / "data"
PKG_DATA = ROOT / "outputs" / "2026_06_14" / "data"
BASE = ROOT / "outputs" / "2026_06_14"
OUT = BASE / "phase_1_to_5_analysis"
OUT_DATA = OUT / "data"
OUT_FIG = OUT / "figures"
OUT_REPORT = OUT / "reports"
PPT_OUT = BASE / "PZONE_bottleneck_phase1_5_summary_7slides.pptx"
SCRIPT_OUT = BASE / "PZONE_bottleneck_phase1_5_ppt_script.md"
SQLITE = ROOT / "data" / "pzone_analysis.sqlite"

FONT = "Malgun Gothic"
NAVY = RGBColor(24, 39, 61)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(185, 28, 28)
GRAY = RGBColor(75, 85, 99)
LIGHT_GRAY = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs() -> None:
    for path in [OUT, OUT_DATA, OUT_FIG, OUT_REPORT]:
        path.mkdir(parents=True, exist_ok=True)


def read_csv(name: str, source: Path = SRC_DATA) -> pd.DataFrame:
    path = source / name
    if not path.exists():
        path = PKG_DATA / name
    return pd.read_csv(path)


def product_family(prd_cd: str) -> str:
    if str(prd_cd) in {"PRD1000", "PRD2000"}:
        return "HANDLE"
    if str(prd_cd) in {"PRD3001", "PRD3002"}:
        return "ROOM_MIRROR"
    return "OTHER"


def clean_by_iqr(df: pd.DataFrame, value_col: str, group_cols: list[str] | None = None, hard_limit: float | None = None) -> pd.Series:
    if df.empty:
        return pd.Series(dtype=bool)
    flag = pd.Series(False, index=df.index)
    valid = pd.to_numeric(df[value_col], errors="coerce").notna()
    if group_cols and len(group_cols) == 1:
        groups = df.loc[valid].groupby(group_cols[0], dropna=False).groups
    elif group_cols:
        groups = df.loc[valid].groupby(group_cols, dropna=False).groups
    else:
        groups = {"_all": df.loc[valid].index}
    for _, idx in groups.items():
        series = pd.to_numeric(df.loc[idx, value_col], errors="coerce").dropna()
        if series.empty:
            continue
        q1 = series.quantile(0.25)
        q3 = series.quantile(0.75)
        upper = q3 + 1.5 * (q3 - q1)
        threshold = upper if hard_limit is None else min(upper, hard_limit)
        flag.loc[idx] = pd.to_numeric(df.loc[idx, value_col], errors="coerce") > threshold
    return flag


def q(series: pd.Series, quantile: float) -> float:
    s = pd.to_numeric(series, errors="coerce").dropna()
    return float(s.quantile(quantile)) if not s.empty else math.nan


def fmt_sec(value: float) -> str:
    if pd.isna(value):
        return "-"
    return f"{float(value):,.1f}초"


def fmt_num(value: float, digits: int = 1) -> str:
    if pd.isna(value):
        return "-"
    return f"{float(value):,.{digits}f}"


def zscore(series: pd.Series) -> pd.Series:
    s = pd.to_numeric(series, errors="coerce").fillna(0)
    std = s.std(ddof=0)
    if std == 0:
        return pd.Series(0.0, index=s.index)
    return (s - s.mean()) / std


def write_md(path: Path, text: str) -> None:
    path.write_text(text.strip() + "\n", encoding="utf-8")


def build_data_validation(data: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    all_routes = data["all_routes"]
    routes = data["routes"]
    excluded = data["excluded_routes"]
    events = data["events"]
    waits = data["wait_events"]
    anomalies = data["pairing_anomalies"]
    normal_route_summary = data["normal_route_summary"]

    all_serials = set(all_routes["PRD_SRL_NO"].dropna())
    normal_serials = set(routes["PRD_SRL_NO"].dropna())
    excluded_serials = set(excluded["PRD_SRL_NO"].dropna()) if not excluded.empty else set()
    paired_count = len(events)
    anomaly_count = len(anomalies)
    total_pair_attempts = paired_count + anomaly_count
    pairing_success_rate = paired_count / total_pair_attempts if total_pair_attempts else math.nan

    quality_rows = [
        ["total_serial_before_filter", len(all_serials), "전체 route 추출 serial 수"],
        ["normal_serial_after_filter", len(normal_serials), "정상 route/prefix 기준 분석 serial 수"],
        ["excluded_serial", len(excluded_serials), "제품별 정상 route에서 벗어난 serial 수"],
        ["paired_processing_events", paired_count, "STR/END 페어링 성공 처리 이벤트 수"],
        ["pairing_anomalies", anomaly_count, "STR/END 페어링 실패 또는 잔여 이벤트 수"],
        ["pairing_success_rate", pairing_success_rate, "paired / (paired + anomaly)"],
        ["negative_processing_events", int(events["negative_duration_flag"].sum()), "음수 처리시간 이벤트 수"],
        ["processing_outliers", int(events["outlier_flag"].sum()), "처리시간 이상치 이벤트 수"],
        ["transition_wait_events", len(waits), "공정 전이 대기 이벤트 수"],
        ["negative_wait_events", int(waits["negative_wait_flag"].sum()), "음수 대기시간 이벤트 수"],
    ]
    data_quality = pd.DataFrame(quality_rows, columns=["metric", "value", "description"])

    validation = normal_route_summary.copy()
    normal_counts = routes.groupby("PRD_CD")["PRD_SRL_NO"].nunique().rename("normal_serial_actual").reset_index()
    before_counts = all_routes.groupby("PRD_CD")["PRD_SRL_NO"].nunique().rename("serial_before_filter").reset_index()
    validation = validation.merge(normal_counts, on="PRD_CD", how="left").merge(before_counts, on="PRD_CD", how="left")
    validation["excluded_serial_actual"] = validation["serial_before_filter"] - validation["normal_serial_actual"]
    validation["actual_normal_ratio"] = validation["normal_serial_actual"] / validation["serial_before_filter"]

    if anomalies.empty:
        anomaly_summary = pd.DataFrame(columns=["type", "count", "unique_serials"])
    else:
        anomaly_summary = (
            anomalies.groupby("type", dropna=False)
            .agg(count=("type", "size"), unique_serials=("PRD_SRL_NO", "nunique"))
            .reset_index()
            .sort_values("count", ascending=False)
        )

    data_quality.to_csv(OUT_DATA / "data_quality_summary.csv", index=False, encoding="utf-8-sig")
    validation.to_csv(OUT_DATA / "route_filter_validation.csv", index=False, encoding="utf-8-sig")
    anomaly_summary.to_csv(OUT_DATA / "pairing_anomaly_summary.csv", index=False, encoding="utf-8-sig")

    write_md(
        OUT_REPORT / "01_data_validation.md",
        f"""
# 01. 데이터 검증

## 핵심 결과
- 전체 serial: {len(all_serials)}개
- 정상 route/prefix serial: {len(normal_serials)}개
- 제외 serial: {len(excluded_serials)}개
- STR/END 페어링 성공률: {pairing_success_rate:.2%}
- 처리시간 이상치 이벤트: {int(events['outlier_flag'].sum())}건
- 음수 대기 이벤트: {int(waits['negative_wait_flag'].sum())}건

## 해석
정상 제품/정상 route 기준으로 분석 대상을 고정했기 때문에 `A5 -> A3` 같은 예외 전이는 본 분석에서 제외된다. 제품군별 정상 route 차이는 유지한다. 예를 들어 룸미러 계열의 `A7 -> A5`는 정상 route이다.
""",
    )
    return {"data_quality": data_quality, "route_validation": validation, "anomaly_summary": anomaly_summary}


def build_a3_analysis(data: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    events = data["events"].copy()
    waits = data["wait_events"].copy()
    wip = data["wip_summary"].copy()

    a3_proc = events[(events["logical_process"] == "A3") & (~events["negative_duration_flag"])].copy()
    a3_proc_clean = a3_proc[~a3_proc["outlier_flag"]].copy()

    a3_waits = waits[(waits["from_process"] == "A3") & (~waits["negative_wait_flag"])].copy()
    a3_waits["transition"] = a3_waits["from_process"] + "->" + a3_waits["to_process"]
    a3_waits["outlier_flag"] = clean_by_iqr(a3_waits, "waiting_sec", ["transition"], hard_limit=3600)
    a3_wait_clean = a3_waits[~a3_waits["outlier_flag"]].copy()

    wait_dist = (
        a3_waits.groupby("transition", dropna=False)
        .agg(count_raw=("waiting_sec", "size"), median_wait_sec=("waiting_sec", "median"), p90_wait_sec=("waiting_sec", lambda s: q(s, 0.9)), max_wait_sec=("waiting_sec", "max"), outlier_count=("outlier_flag", "sum"))
        .reset_index()
    )
    clean_dist = (
        a3_wait_clean.groupby("transition", dropna=False)
        .agg(count_clean=("waiting_sec", "size"), median_wait_sec_clean=("waiting_sec", "median"), p90_wait_sec_clean=("waiting_sec", lambda s: q(s, 0.9)), max_wait_sec_clean=("waiting_sec", "max"))
        .reset_index()
    )
    wait_dist = wait_dist.merge(clean_dist, on="transition", how="left").sort_values("p90_wait_sec_clean", ascending=False)

    def wait_p90(transition: str) -> float:
        row = wait_dist.loc[wait_dist["transition"] == transition]
        return float(row["p90_wait_sec_clean"].iloc[0]) if not row.empty and pd.notna(row["p90_wait_sec_clean"].iloc[0]) else math.nan

    buffer_wip = wip.loc[wip["area"] == "BUFFER", "p90_wip"]
    finished_wip = wip.loc[wip["area"] == "FINISHED_STORAGE", "p90_wip"]
    decomposition = pd.DataFrame(
        [
            {
                "component": "A3_PROCESS",
                "metric_value": q(a3_proc_clean["processing_sec"], 0.9),
                "unit": "sec",
                "evidence": "A3 clean processing p90",
                "interpretation": "A3 자체 처리시간은 후단 대기 대비 작다.",
                "certainty": "High",
            },
            {
                "component": "AMR_DISPATCH_WAIT",
                "metric_value": wait_p90("A3->AMR_LD250"),
                "unit": "sec",
                "evidence": "A3 -> AMR_LD250 clean waiting p90",
                "interpretation": "AMR 또는 후단 반출 대기 후보. 호출/배정 로그가 없어 배정 지연 확정은 불가.",
                "certainty": "Medium",
            },
            {
                "component": "C2_STORAGE_WAIT",
                "metric_value": wait_p90("A3->C2_FINISHED_STORAGE"),
                "unit": "sec",
                "evidence": "A3 -> C2 clean waiting p90",
                "interpretation": "완제품 창고 적재/반출 대기 후보. C2 slot capacity 로그가 없어 포화 확정은 불가.",
                "certainty": "Medium",
            },
            {
                "component": "BUFFER_ACCUMULATION",
                "metric_value": float(buffer_wip.iloc[0]) if not buffer_wip.empty else math.nan,
                "unit": "wip",
                "evidence": "BUFFER WIP p90",
                "interpretation": "후단 배출 지연이 버퍼 누적으로 나타난다.",
                "certainty": "High",
            },
            {
                "component": "FINISHED_STORAGE_ACCUMULATION",
                "metric_value": float(finished_wip.iloc[0]) if not finished_wip.empty else math.nan,
                "unit": "wip",
                "evidence": "FINISHED_STORAGE WIP p90",
                "interpretation": "현재 로그상 완제품 창고 WIP 자체는 낮게 관측된다.",
                "certainty": "Medium",
            },
        ]
    )

    decomposition.to_csv(OUT_DATA / "a3_bottleneck_decomposition.csv", index=False, encoding="utf-8-sig")
    wait_dist.to_csv(OUT_DATA / "a3_wait_distribution.csv", index=False, encoding="utf-8-sig")

    plot_df = decomposition.copy()
    plot_df["plot_value"] = plot_df["metric_value"]
    plt.figure(figsize=(9, 4.6))
    sns.barplot(data=plot_df, y="component", x="plot_value", hue="component", dodge=False, legend=False, palette="viridis")
    plt.xlabel("Value (seconds for time components, WIP for WIP components)")
    plt.ylabel("")
    plt.title("A3 Bottleneck Decomposition")
    plt.tight_layout()
    plt.savefig(OUT_FIG / "a3_decomposition.png", dpi=160)
    plt.close()

    write_md(
        OUT_REPORT / "02_a3_deep_dive.md",
        f"""
# 02. A3 병목 세분화

## 핵심 결과
- A3 처리 p90(clean): {fmt_sec(decomposition.loc[decomposition['component'] == 'A3_PROCESS', 'metric_value'].iloc[0])}
- A3 -> AMR_LD250 대기 p90(clean): {fmt_sec(decomposition.loc[decomposition['component'] == 'AMR_DISPATCH_WAIT', 'metric_value'].iloc[0])}
- A3 -> C2 대기 p90(clean): {fmt_sec(decomposition.loc[decomposition['component'] == 'C2_STORAGE_WAIT', 'metric_value'].iloc[0])}
- BUFFER WIP p90: {fmt_num(decomposition.loc[decomposition['component'] == 'BUFFER_ACCUMULATION', 'metric_value'].iloc[0])}

## 판단
A3는 설비 처리시간 자체보다 A3 이후 AMR/C2 반출 대기와 버퍼 누적이 더 큰 병목 신호다. 단, AMR 호출/배정 로그와 C2 slot capacity 로그가 없어 AMR 배정 지연 또는 C2 포화를 확정하지는 않는다.
""",
    )
    return {"a3_decomposition": decomposition, "a3_wait_distribution": wait_dist}


def build_a4_analysis(data: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    events = data["events"].copy()
    rail_windows_existing = data["a4_windows"].copy()
    downstream = data["a4_downstream"].copy()
    waits = data["wait_events"].copy()

    a4 = events[(events["logical_process"] == "A4") & (~events["negative_duration_flag"])].copy()
    a4["product_family"] = a4["PRD_CD"].map(product_family)
    threshold = float(a4["processing_sec"].quantile(0.75))
    if not rail_windows_existing.empty and "a4_long_threshold_sec" in rail_windows_existing.columns:
        threshold = float(rail_windows_existing["a4_long_threshold_sec"].dropna().iloc[0])
    a4["window_type"] = a4["processing_sec"].apply(lambda x: "long" if x > threshold else "normal")

    existing_long_cols = [
        "PRD_SRL_NO",
        "rail_event_count_overlap",
        "rail_product_count_overlap",
        "active_rail_count_overlap",
        "same_family_rail_events",
        "other_family_rail_events",
        "other_family_product_count",
        "overlap_wait_event_count",
        "overlap_wait_p90_sec",
    ]
    a4 = a4.merge(rail_windows_existing[[c for c in existing_long_cols if c in rail_windows_existing.columns]], on="PRD_SRL_NO", how="left")
    overlap_cols = [c for c in existing_long_cols if c != "PRD_SRL_NO"]
    for col in overlap_cols:
        if col in a4.columns:
            a4[col] = a4[col].fillna(0)

    long_vs_normal = (
        a4.groupby("window_type")
        .agg(
            window_count=("PRD_SRL_NO", "nunique"),
            processing_p50_sec=("processing_sec", "median"),
            processing_p90_sec=("processing_sec", lambda s: q(s, 0.9)),
            rail_events_overlap_mean=("rail_event_count_overlap", "mean"),
            other_family_events_mean=("other_family_rail_events", "mean"),
            overlap_wait_p90_mean=("overlap_wait_p90_sec", "mean"),
        )
        .reset_index()
    )
    long_row = long_vs_normal.loc[long_vs_normal["window_type"] == "long"]
    normal_row = long_vs_normal.loc[long_vs_normal["window_type"] == "normal"]
    if not long_row.empty and not normal_row.empty:
        ratio = (float(long_row["other_family_events_mean"].iloc[0]) + 1) / (float(normal_row["other_family_events_mean"].iloc[0]) + 1)
    else:
        ratio = math.nan
    evidence_strength = "Strong" if ratio >= 2.0 else "Moderate" if ratio >= 1.2 else "Weak"

    cross_family = pd.DataFrame()
    if not rail_windows_existing.empty:
        cross_family = (
            rail_windows_existing.groupby(["product_family"], dropna=False)
            .agg(
                long_window_count=("PRD_SRL_NO", "nunique"),
                same_family_rail_events=("same_family_rail_events", "sum"),
                other_family_rail_events=("other_family_rail_events", "sum"),
                other_family_product_count_p90=("other_family_product_count", lambda s: q(s, 0.9)),
                overlap_wait_p90_sec_mean=("overlap_wait_p90_sec", "mean"),
            )
            .reset_index()
        )
        cross_family["other_family_ratio"] = cross_family["other_family_rail_events"] / (
            cross_family["same_family_rail_events"] + cross_family["other_family_rail_events"]
        ).replace(0, math.nan)

    downstream_summary = pd.DataFrame()
    if not downstream.empty:
        downstream_summary = (
            downstream.groupby("lag_minutes")
            .agg(
                rail_events_after_mean=("rail_event_count_after", "mean"),
                active_rail_count_after_mean=("active_rail_count_after", "mean"),
                waiting_event_count_after_mean=("waiting_event_count_after", "mean"),
                waiting_p90_sec_after_mean=("waiting_p90_sec_after", "mean"),
            )
            .reset_index()
        )

    evidence = pd.DataFrame(
        [
            ["a4_long_threshold_sec", threshold, "A4 long window 기준"],
            ["a4_long_window_count", int((a4["window_type"] == "long").sum()), "A4 long window 수"],
            ["a4_other_family_overlap_ratio_long_vs_normal", ratio, "long window의 다른 제품군 overlap 평균 / normal 평균"],
            ["evidence_strength", evidence_strength, "A4 blocking 가설 근거 강도"],
        ],
        columns=["metric", "value", "description"],
    )

    long_vs_normal.to_csv(OUT_DATA / "a4_long_vs_normal_comparison.csv", index=False, encoding="utf-8-sig")
    cross_family.to_csv(OUT_DATA / "a4_cross_family_rail_overlap.csv", index=False, encoding="utf-8-sig")
    evidence.to_csv(OUT_DATA / "a4_blocking_evidence_summary.csv", index=False, encoding="utf-8-sig")
    downstream_summary.to_csv(OUT_DATA / "a4_downstream_lag_summary.csv", index=False, encoding="utf-8-sig")

    plt.figure(figsize=(8, 4.6))
    plot_df = long_vs_normal.melt(id_vars=["window_type"], value_vars=["processing_p90_sec", "other_family_events_mean", "overlap_wait_p90_mean"], var_name="metric", value_name="value")
    sns.barplot(data=plot_df, x="metric", y="value", hue="window_type")
    plt.title("A4 Long vs Normal Window")
    plt.xlabel("")
    plt.ylabel("Value")
    plt.xticks(rotation=18, ha="right")
    plt.tight_layout()
    plt.savefig(OUT_FIG / "a4_long_vs_normal.png", dpi=160)
    plt.close()

    if not cross_family.empty:
        plt.figure(figsize=(8, 4.6))
        cf = cross_family.melt(id_vars=["product_family"], value_vars=["same_family_rail_events", "other_family_rail_events"], var_name="overlap_type", value_name="rail_events")
        sns.barplot(data=cf, x="product_family", y="rail_events", hue="overlap_type")
        plt.title("A4 Long Window Rail Overlap by Product Family")
        plt.xlabel("")
        plt.ylabel("Rail events")
        plt.tight_layout()
        plt.savefig(OUT_FIG / "a4_cross_family_overlap.png", dpi=160)
        plt.close()

    write_md(
        OUT_REPORT / "03_a4_shared_rail_validation.md",
        f"""
# 03. A4 공유 레일 Blocking 검증

## 핵심 결과
- A4 long window 기준: {fmt_sec(threshold)}
- A4 long window 수: {int((a4['window_type'] == 'long').sum())}개
- long/normal 다른 제품군 overlap 비율: {fmt_num(ratio, 2)}
- A4 blocking 가설 근거 강도: **{evidence_strength}**

## 판단
A4는 처리시간이 긴 복합 CELL이며, long window 중 공유 레일 이벤트와 다른 제품군 overlap이 관측된다. 다만 `TRN_DEV_ID`의 실제 물리 위치가 DB에 없으므로, 현재 결론은 공유 레일 blocking 후보 검증 결과로 해석해야 한다.
""",
    )
    return {"a4_long_vs_normal": long_vs_normal, "a4_cross_family": cross_family, "a4_evidence": evidence}


def build_product_family_analysis(data: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    events = data["events"].copy()
    waits = data["wait_events"].copy()
    wip = data["wip_summary"].copy()
    occ = data["occupancy"].copy()
    oee = data["oee"].copy()
    routes = data["normal_route_summary"].copy()

    events["product_family"] = events["PRD_CD"].map(product_family)
    waits["product_family"] = waits["PRD_CD"].map(product_family)

    route_summary = routes.copy()
    route_summary["product_family"] = route_summary["PRD_CD"].map(product_family)

    proc = (
        events[(~events["negative_duration_flag"]) & (~events["outlier_flag"])]
        .groupby(["product_family", "logical_process"], dropna=False)
        .agg(processing_count=("processing_sec", "size"), processing_p90_sec=("processing_sec", lambda s: q(s, 0.9)))
        .reset_index()
    )
    waits_valid = waits[~waits["negative_wait_flag"]].copy()
    waits_valid["transition"] = waits_valid["from_process"] + "->" + waits_valid["to_process"]
    waits_valid["outlier_flag"] = clean_by_iqr(waits_valid, "waiting_sec", ["product_family", "transition"], hard_limit=3600)
    waits_clean = waits_valid[~waits_valid["outlier_flag"]]
    transition = (
        waits_clean.groupby(["product_family", "from_process", "to_process", "transition"], dropna=False)
        .agg(wait_count=("waiting_sec", "size"), waiting_p90_sec=("waiting_sec", lambda s: q(s, 0.9)), waiting_median_sec=("waiting_sec", "median"))
        .reset_index()
        .sort_values(["product_family", "waiting_p90_sec"], ascending=[True, False])
    )

    wait_for_score = (
        transition.groupby(["product_family", "from_process"], dropna=False)["waiting_p90_sec"].max().rename("wait_from_p90").reset_index().rename(columns={"from_process": "logical_process"})
    )
    wait_to = (
        transition.groupby(["product_family", "to_process"], dropna=False)["waiting_p90_sec"].max().rename("wait_to_p90").reset_index().rename(columns={"to_process": "logical_process"})
    )
    score = proc.merge(wait_for_score, on=["product_family", "logical_process"], how="outer").merge(wait_to, on=["product_family", "logical_process"], how="outer")
    score["waiting_p90_sec"] = score[["wait_from_p90", "wait_to_p90"]].max(axis=1).fillna(0)
    score["processing_p90_sec"] = score["processing_p90_sec"].fillna(0)
    score["wip_p90"] = 0.0
    score.loc[score["logical_process"] == "A3", "wip_p90"] = float(wip.loc[wip["area"] == "BUFFER", "p90_wip"].iloc[0]) if not wip.loc[wip["area"] == "BUFFER"].empty else 0.0
    score.loc[score["logical_process"] == "C2_FINISHED_STORAGE", "wip_p90"] = float(wip.loc[wip["area"] == "FINISHED_STORAGE", "p90_wip"].iloc[0]) if not wip.loc[wip["area"] == "FINISHED_STORAGE"].empty else 0.0

    occ_proc = occ.groupby("logical_process")["p90_occupancy_sec_clean"].max().rename("occupancy_p90_sec").reset_index() if not occ.empty else pd.DataFrame(columns=["logical_process", "occupancy_p90_sec"])
    alarm = oee[["logical_process", "alarm_ratio"]].copy() if not oee.empty else pd.DataFrame(columns=["logical_process", "alarm_ratio"])
    score = score.merge(occ_proc, on="logical_process", how="left").merge(alarm, on="logical_process", how="left")
    score[["occupancy_p90_sec", "alarm_ratio"]] = score[["occupancy_p90_sec", "alarm_ratio"]].fillna(0)
    score["bottleneck_score"] = 0.0
    parts = []
    for fam, g in score.groupby("product_family", dropna=False):
        tmp = g.copy()
        tmp["bottleneck_score"] = (
            zscore(tmp["processing_p90_sec"])
            + zscore(tmp["waiting_p90_sec"])
            + zscore(tmp["wip_p90"])
            + zscore(tmp["occupancy_p90_sec"])
            + zscore(tmp["alarm_ratio"])
        )
        tmp = tmp.sort_values("bottleneck_score", ascending=False).reset_index(drop=True)
        tmp.insert(0, "family_rank", range(1, len(tmp) + 1))
        parts.append(tmp)
    ranking = pd.concat(parts, ignore_index=True) if parts else pd.DataFrame()
    ranking = ranking.rename(columns={"logical_process": "process"})

    route_summary.to_csv(OUT_DATA / "product_family_route_summary.csv", index=False, encoding="utf-8-sig")
    ranking.to_csv(OUT_DATA / "product_family_bottleneck_ranking.csv", index=False, encoding="utf-8-sig")
    transition.to_csv(OUT_DATA / "product_family_transition_waiting.csv", index=False, encoding="utf-8-sig")

    top = ranking[ranking["family_rank"] <= 6].copy()
    plt.figure(figsize=(9, 5))
    sns.barplot(data=top, y="process", x="bottleneck_score", hue="product_family")
    plt.title("Bottleneck Ranking by Product Family")
    plt.xlabel("Bottleneck score")
    plt.ylabel("")
    plt.tight_layout()
    plt.savefig(OUT_FIG / "product_family_bottleneck_comparison.png", dpi=160)
    plt.close()

    handle_top = ranking[ranking["product_family"] == "HANDLE"].head(3)["process"].tolist()
    mirror_top = ranking[ranking["product_family"] == "ROOM_MIRROR"].head(3)["process"].tolist()
    write_md(
        OUT_REPORT / "04_product_family_analysis.md",
        f"""
# 04. 제품군별 병목 분리

## 핵심 결과
- HANDLE 계열 상위 병목 후보: {', '.join(handle_top)}
- ROOM_MIRROR 계열 상위 병목 후보: {', '.join(mirror_top)}

## 판단
핸들 계열은 A4/A6/A9_2를 포함하는 긴 route를 가지며, 룸미러 계열은 A7 -> A5 -> A9_1 -> A3 흐름이 정상이다. 따라서 제품군을 합쳐 보면 route 차이와 병목 원인이 섞일 수 있으므로 제품군별 병목 순위와 스케줄링 규칙을 별도로 관리해야 한다.
""",
    )
    return {"family_routes": route_summary, "family_ranking": ranking, "family_transition": transition}


def load_transport_events(normal_serials: set[str]) -> pd.DataFrame:
    conn = sqlite3.connect(SQLITE)
    df = pd.read_sql_query(
        """
        SELECT PRD_SRL_NO, PRD_CD, CMP_EQ_ID, TRN_CD, TRN_DEV_ID, MBL_NMBR, PRD_WRK_CD, TRN_QNT, REG_DT
        FROM prc_trns_tb
        WHERE CMP_CD='4188300219' AND CMP_LINE_ID='LN01' AND PRD_TYP_NO='PRDC'
          AND PRD_CD IN ('PRD1000','PRD2000','PRD3001','PRD3002')
        """,
        conn,
    )
    conn.close()
    df = df[df["PRD_SRL_NO"].isin(normal_serials)].copy()
    df["REG_DT"] = pd.to_datetime(df["REG_DT"], errors="coerce")
    df["TRN_DEV_ID"] = pd.to_numeric(df["TRN_DEV_ID"], errors="coerce")
    df["TRN_QNT"] = pd.to_numeric(df["TRN_QNT"], errors="coerce")
    df["product_family"] = df["PRD_CD"].map(product_family)
    return df


def build_transport_analysis(data: dict[str, pd.DataFrame]) -> dict[str, pd.DataFrame]:
    events = data["events"].copy()
    normal_serials = set(events["PRD_SRL_NO"].dropna())
    trn = load_transport_events(normal_serials)
    rail = trn[(trn["TRN_CD"] == "TR01") & (trn["TRN_DEV_ID"].between(1, 4))].copy()
    rail["bucket_5min"] = rail["REG_DT"].dt.floor("5min")

    total_buckets = max(1, rail["bucket_5min"].nunique())
    utilization = (
        rail.groupby("TRN_DEV_ID", dropna=False)
        .agg(
            event_count=("PRD_SRL_NO", "size"),
            product_count=("PRD_SRL_NO", "nunique"),
            active_bucket_count=("bucket_5min", "nunique"),
            trn_qnt_sum=("TRN_QNT", "sum"),
            trn_qnt_mean=("TRN_QNT", "mean"),
        )
        .reset_index()
    )
    utilization["active_bucket_ratio"] = utilization["active_bucket_count"] / total_buckets

    occ = data["occupancy"].copy()
    rail_occ = occ[(occ["TRN_CD"] == "TR01") & (pd.to_numeric(occ["TRN_DEV_ID"], errors="coerce").between(1, 4))].copy()
    rail_occ["TRN_DEV_ID"] = pd.to_numeric(rail_occ["TRN_DEV_ID"], errors="coerce")
    occ_by_dev = (
        rail_occ.groupby("TRN_DEV_ID")
        .agg(p90_occupancy_sec_clean=("p90_occupancy_sec_clean", "max"), median_occupancy_sec_clean=("median_occupancy_sec_clean", "median"))
        .reset_index()
    )
    utilization = utilization.merge(occ_by_dev, on="TRN_DEV_ID", how="left").sort_values("TRN_DEV_ID")

    family_usage = (
        rail.groupby(["TRN_DEV_ID", "product_family"], dropna=False)
        .agg(event_count=("PRD_SRL_NO", "size"), product_count=("PRD_SRL_NO", "nunique"))
        .reset_index()
        .sort_values(["TRN_DEV_ID", "product_family"])
    )

    bottleneck = data["bottleneck"].copy()
    waits = data["wait_summary"].copy()
    oee = data["oee"].copy()
    amr_rows = []
    for process in ["AMR_LD90", "AMR_LD250"]:
        b = bottleneck[bottleneck["process"] == process]
        o = oee[oee["logical_process"] == process]
        related = waits[(waits["from_process"] == process) | (waits["to_process"] == process)]
        amr_rows.append(
            {
                "resource": process,
                "bottleneck_rank": int(b["rank"].iloc[0]) if not b.empty else math.nan,
                "processing_p90_sec": float(b["processing_p90_sec"].iloc[0]) if not b.empty else math.nan,
                "waiting_p90_sec": float(b["waiting_p90_sec"].iloc[0]) if not b.empty else math.nan,
                "alarm_ratio": float(o["alarm_ratio"].iloc[0]) if not o.empty else math.nan,
                "run_ratio": float(o["run_ratio"].iloc[0]) if not o.empty else math.nan,
                "related_transition_count": int(len(related)),
                "top_related_transition": related.sort_values("p90_wait_sec_clean", ascending=False)["transition"].iloc[0] if not related.empty else "",
                "top_related_wait_p90_sec": float(related.sort_values("p90_wait_sec_clean", ascending=False)["p90_wait_sec_clean"].iloc[0]) if not related.empty else math.nan,
            }
        )
    amr_summary = pd.DataFrame(amr_rows)

    limitation = pd.DataFrame(
        [
            ["TRN_DEV_ID_PHYSICAL_LOCATION", "Missing", "TRN_DEV_ID=1~4의 실제 물리 위치/구간 매핑 필요"],
            ["AMR_DISPATCH_TIMESTAMPS", "Missing", "AMR 호출, 배정, 취소, 대기 상태 로그 필요"],
            ["C2_SLOT_CAPACITY", "Missing", "C2 완제품 창고 slot capacity와 포화 상태 로그 필요"],
            ["OEE_IDLE_DEFINITION", "Needs confirmation", "IDLE이 정상 대기인지 설비 비가동인지 코드 정의 필요"],
            ["BUFFER_STOPPER_CAPACITY", "Missing", "스토퍼/버퍼별 실제 capacity와 blocked 상태 로그 필요"],
        ],
        columns=["information", "status", "why_needed"],
    )

    utilization.to_csv(OUT_DATA / "rail_device_utilization.csv", index=False, encoding="utf-8-sig")
    family_usage.to_csv(OUT_DATA / "rail_device_product_family_usage.csv", index=False, encoding="utf-8-sig")
    amr_summary.to_csv(OUT_DATA / "amr_resource_summary.csv", index=False, encoding="utf-8-sig")
    limitation.to_csv(OUT_DATA / "transport_resource_limitation_summary.csv", index=False, encoding="utf-8-sig")

    plt.figure(figsize=(8, 4.6))
    sns.barplot(data=utilization, x="TRN_DEV_ID", y="active_bucket_ratio", hue="TRN_DEV_ID", dodge=False, legend=False, palette="Blues_d")
    plt.title("Rail Device Active Bucket Ratio")
    plt.xlabel("TRN_DEV_ID")
    plt.ylabel("Active bucket ratio")
    plt.tight_layout()
    plt.savefig(OUT_FIG / "rail_device_utilization.png", dpi=160)
    plt.close()

    plt.figure(figsize=(7, 4.4))
    sns.barplot(data=amr_summary, x="resource", y="top_related_wait_p90_sec", hue="resource", dodge=False, legend=False, palette="Oranges")
    plt.title("AMR Top Related Transition Wait p90")
    plt.xlabel("")
    plt.ylabel("Seconds")
    plt.tight_layout()
    plt.savefig(OUT_FIG / "amr_wait_summary.png", dpi=160)
    plt.close()

    write_md(
        OUT_REPORT / "05_rail_amr_analysis.md",
        f"""
# 05. 레일/AMR 이송 자원 분석

## 핵심 결과
- TRN_DEV_ID=1~4별 이벤트 수와 active bucket ratio를 계산했다.
- AMR_LD90과 AMR_LD250을 분리해 병목 순위, 대기 p90, alarm ratio를 요약했다.
- 현재 DB만으로는 TRN_DEV_ID의 물리 위치와 AMR 배정 지연 원인을 확정할 수 없다.

## 판단
레일/AMR은 공정 처리시간보다 공정 간 대기를 증폭시키는 공유 이송 자원으로 봐야 한다. 다만 물리 위치 매핑과 AMR dispatch 로그가 없으므로 현재 결과는 자원 병목 후보를 좁히는 분석으로 사용한다.
""",
    )
    return {"rail_utilization": utilization, "rail_family_usage": family_usage, "amr_summary": amr_summary, "limitations": limitation}


def write_required_gap_report(limitations: pd.DataFrame) -> None:
    rows = "\n".join(f"- **{r.information}**: {r.why_needed}" for r in limitations.itertuples())
    write_md(
        OUT / "required_field_gap_report.md",
        f"""
# Required Field Gap Report

현재 DB/CSV 로그만으로 원인을 확정하기 어려운 정보는 아래와 같다.

{rows}

## 사용 원칙
위 정보가 없을 때는 AMR 지연, C2 포화, 레일 물리 blocking을 확정 표현하지 않는다. 대신 로그 기반 후보 원인으로 분류하고, 현장 매핑이 확보되면 재검증한다.
""",
    )


def make_markdown_summary(results: dict[str, pd.DataFrame]) -> None:
    dq = results["validation"]["data_quality"]
    a3 = results["a3"]["a3_decomposition"]
    a4 = results["a4"]["a4_evidence"]
    fam = results["family"]["family_ranking"]
    amr = results["transport"]["amr_summary"]
    data_quality = {r.metric: r.value for r in dq.itertuples()}
    a4_strength = a4.loc[a4["metric"] == "evidence_strength", "value"].iloc[0]
    handle_top = ", ".join(fam[fam["product_family"] == "HANDLE"].head(3)["process"].tolist())
    mirror_top = ", ".join(fam[fam["product_family"] == "ROOM_MIRROR"].head(3)["process"].tolist())
    write_md(
        OUT / "phase_1_to_5_summary.md",
        f"""
# P-ZONE 1~5번 분석 요약

## 완료 범위
- 1. 데이터/분석 검증
- 2. A3 병목 세분화
- 3. A4 공유 레일 blocking 검증
- 4. 제품군별 병목 분리
- 5. 레일/AMR 이송 자원 분석 고도화

## 핵심 결론
- 정상 route 분석 serial: {int(float(data_quality['normal_serial_after_filter']))} / {int(float(data_quality['total_serial_before_filter']))}
- A3는 자체 처리보다 AMR/C2 반출 대기와 BUFFER WIP가 큰 후단 배출 병목이다.
- A4는 long window와 공유 레일 overlap이 관측되는 blocking 후보이며, 현재 근거 강도는 **{a4_strength}**이다.
- HANDLE 상위 병목 후보: {handle_top}
- ROOM_MIRROR 상위 병목 후보: {mirror_top}
- 레일/AMR 분석은 가능하지만 TRN_DEV_ID 물리 위치와 AMR dispatch 로그가 없어 원인 확정에는 추가 정보가 필요하다.

## 개선 단계 진입 조건
1~5번 분석으로 병목 후보와 원인 가설은 정리되었다. 다음 단계는 rule-based scheduling baseline을 설계하고 로그 리플레이/시뮬레이션으로 효과를 검증하는 것이다.
""",
    )


def set_text_style(shape, font_size=14, color=NAVY, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, font_size=14, color=NAVY, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    set_text_style(shape, font_size, color, bold)
    return shape


def add_title(slide, title, subtitle=None, section=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_textbox(slide, 0.45, 0.35, 10.4, 0.45, title, 23, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.48, 0.82, 10.8, 0.32, subtitle, 10.5, GRAY)
    if section:
        add_textbox(slide, 11.65, 0.42, 1.2, 0.3, section, 9.5, GRAY, False, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, font_size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(4)
    return box


def add_callout(slide, x, y, w, h, title, body, color=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(239, 246, 255)
    rect.line.color.rgb = color
    rect.line.width = Pt(1)
    add_textbox(slide, x + 0.15, y + 0.1, w - 0.3, 0.24, title, 10, color, True)
    add_textbox(slide, x + 0.15, y + 0.38, w - 0.3, h - 0.48, body, 11, NAVY)


def add_image_fit(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio >= box_ratio:
        final_w = w
        final_h = w / ratio
    else:
        final_h = h
        final_w = h * ratio
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))


def add_table(slide, x, y, w, h, headers, rows, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.035)
            cell.margin_right = Inches(0.035)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = NAVY
    for cell in table.rows[0].cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
    return table_shape


def add_footer(slide, page):
    add_textbox(slide, 0.45, 7.12, 5.2, 0.2, "P-ZONE 1~5번 병목 진단 | 로그 기반 분석", 8.5, GRAY)
    add_textbox(slide, 12.3, 7.12, 0.55, 0.2, f"{page}/7", 8.5, GRAY, False, PP_ALIGN.RIGHT)


def add_notes(slide, note: str):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = note
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(10)


def build_ppt(results: dict[str, dict[str, pd.DataFrame]]) -> dict[int, str]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes: dict[int, str] = {}

    dq = results["validation"]["data_quality"]
    dq_map = {r.metric: r.value for r in dq.itertuples()}
    bottleneck = results["base"]["bottleneck"]
    a3 = results["a3"]["a3_decomposition"]
    a3_wait = results["a3"]["a3_wait_distribution"]
    a4_evidence = results["a4"]["a4_evidence"]
    a4_strength = a4_evidence.loc[a4_evidence["metric"] == "evidence_strength", "value"].iloc[0]
    family_rank = results["family"]["family_ranking"]
    rail_util = results["transport"]["rail_utilization"]
    amr = results["transport"]["amr_summary"]

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "분석 목적과 현재 완료 범위", "1~5번은 병목 원인 진단, 6번부터 개선 규칙 설계 단계", "01")
    add_callout(slide, 0.6, 1.25, 5.7, 1.25, "목적", "P-ZONE 전체 공정에서 병목 위치와 원인을 찾고, 스케줄링으로 개선 가능한 구간과 물리 제약 구간을 구분한다.", BLUE)
    add_table(
        slide,
        0.65,
        2.85,
        5.8,
        2.05,
        ["완료 단계", "내용"],
        [
            ["1", "데이터/route/페어링 검증"],
            ["2", "A3 후단 병목 세분화"],
            ["3", "A4 공유 레일 blocking 검증"],
            ["4", "제품군별 병목 분리"],
            ["5", "레일/AMR 자원 분석"],
        ],
        9.5,
    )
    add_image_fit(slide, BASE / "figures" / "01_bottleneck_ranking.png", 6.8, 1.15, 5.6, 4.15)
    add_callout(slide, 0.7, 5.65, 11.7, 0.65, "핵심 결론", "A3는 후단 반출 병목, A4는 공유 레일 blocking 후보, 제품군별 route 차이와 레일/AMR 제약은 병목을 증폭시키는 연결 요인이다.", ORANGE)
    notes[1] = "이번 분석은 전체 공정을 같은 기준으로 먼저 본 뒤, 근거가 강한 병목을 세부적으로 나누는 방식입니다. 1번부터 5번까지는 진단 단계이고, 6번부터가 실제 개선 규칙을 설계하는 단계입니다. 현재 데이터 기준으로 가장 강한 병목은 A3 후단 반출이며, A4는 공유 레일 blocking 후보로 확인했습니다."
    add_notes(slide, notes[1])
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "데이터 검증과 정상 Route 기준", "예외 route를 제거하고 제품군별 정상 흐름만 분석", "02")
    add_table(
        slide,
        0.55,
        1.25,
        5.4,
        2.1,
        ["지표", "값"],
        [
            ["전체 serial", int(float(dq_map["total_serial_before_filter"]))],
            ["정상 route serial", int(float(dq_map["normal_serial_after_filter"]))],
            ["제외 serial", int(float(dq_map["excluded_serial"]))],
            ["페어링 성공률", f"{float(dq_map['pairing_success_rate']):.2%}"],
            ["처리 이상치", int(float(dq_map["processing_outliers"]))],
        ],
        10,
    )
    route_rows = []
    for _, r in results["validation"]["route_validation"].iterrows():
        route = str(r["normal_route"]).replace("C1_RAW_STORAGE", "C1").replace("C2_FINISHED_STORAGE", "C2")
        if len(route) > 74:
            route = route[:71] + "..."
        route_rows.append([r["PRD_CD"], route, int(r["normal_serial_actual"])])
    add_table(slide, 6.25, 1.25, 6.6, 2.3, ["제품", "정상 route", "분석 serial"], route_rows, 8.1)
    add_bullets(
        slide,
        0.75,
        3.95,
        11.7,
        1.35,
        [
            "A5 -> A3 같은 비정상 전이는 제외했다.",
            "룸미러 계열의 A7 -> A5는 정상 route이므로 유지했다.",
            "clean 기준은 정상 route/prefix + 음수/누락/극단 이상치 분리 기준이다.",
        ],
        12.3,
    )
    add_callout(slide, 0.7, 5.75, 11.7, 0.55, "의미", "정상 route 기준을 먼저 고정해야 제품군별 정상 차이와 비정상 로그가 섞여 병목이 왜곡되는 것을 막을 수 있다.", TEAL)
    notes[2] = "두 번째 장은 분석 신뢰도를 설명하는 장입니다. 전체 324개 serial 중 정상 route 또는 정상 prefix에 해당하는 294개를 분석했습니다. 여기서 중요한 점은 제품군별 route 차이입니다. A7에서 A5로 가는 흐름은 룸미러 계열에서는 정상이고, A5에서 A3로 바로 가는 흐름은 정상 route에서 벗어나 제외했습니다."
    add_notes(slide, notes[2])
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "전체 병목 스크리닝 결과", "처리시간, 대기시간, WIP, 점유시간, alarm ratio를 표준화해 합산", "03")
    add_image_fit(slide, BASE / "figures" / "01_bottleneck_ranking.png", 0.55, 1.15, 7.1, 4.7)
    rows = []
    for _, r in bottleneck.head(7).iterrows():
        rows.append([int(r["rank"]), r["process"], f"{r['bottleneck_score']:.2f}", fmt_sec(r["processing_p90_sec"]), fmt_sec(r["waiting_p90_sec"])])
    add_table(slide, 7.95, 1.25, 4.9, 2.25, ["순위", "공정", "점수", "처리 p90", "대기 p90"], rows, 7.8)
    add_bullets(
        slide,
        8.05,
        3.85,
        4.55,
        1.65,
        [
            "A3는 대기/WIP 영향으로 1순위다.",
            "A4는 처리시간/점유시간이 높다.",
            "A9/A7/A5/AMR은 연결 병목 후보로 남는다.",
        ],
        11.4,
    )
    add_callout(slide, 7.95, 5.75, 4.9, 0.58, "점수 해석", "절대 점수가 아니라 전체 공정 대비 상대적으로 튀는 병목 지표의 합산값이다.", BLUE)
    notes[3] = "병목 점수는 처리시간, 대기시간, WIP, 점유시간, alarm ratio를 각각 표준화해서 같은 비중으로 더한 상대 점수입니다. A3는 처리시간 자체가 길어서가 아니라 후단 대기와 WIP가 크기 때문에 높게 나왔습니다. A4는 처리시간과 점유시간이 높아 구조적 병목 후보로 분류됩니다."
    add_notes(slide, notes[3])
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "A3 병목 세분화", "A3 자체보다 AMR/C2 반출 대기와 BUFFER 누적이 핵심", "04")
    add_image_fit(slide, OUT_FIG / "a3_decomposition.png", 0.55, 1.18, 6.4, 4.25)
    a3_rows = [[r.component, fmt_num(r.metric_value), r.unit, r.certainty] for r in a3.itertuples()]
    add_table(slide, 7.2, 1.2, 5.45, 2.2, ["요소", "값", "단위", "확실성"], a3_rows, 7.6)
    top_wait_rows = []
    for _, r in a3_wait.head(3).iterrows():
        top_wait_rows.append([r["transition"], fmt_sec(r["p90_wait_sec_clean"]), int(r["count_clean"]) if pd.notna(r["count_clean"]) else 0])
    add_table(slide, 7.2, 3.75, 5.45, 1.15, ["전이", "대기 p90", "clean 수"], top_wait_rows, 8.2)
    add_callout(slide, 0.75, 5.72, 11.7, 0.62, "판단", "A3는 Scheduling Feasible 후보지만, AMR 배정 로그와 C2 slot capacity가 없어 AMR/C2 중 최종 원인은 추가 확인이 필요하다.", ORANGE)
    notes[4] = "A3를 세분화하면 A3 처리 p90은 작지만 A3에서 AMR_LD250 또는 C2로 넘어가는 대기 p90이 훨씬 큽니다. 따라서 A3 병목은 설비 처리 자체보다는 후단 반출 병목으로 보는 것이 맞습니다. 다만 현재 DB에는 AMR 호출과 배정 시각, C2 실제 포화 정보가 없어서 AMR 문제인지 C2 문제인지는 후보 수준으로 구분합니다."
    add_notes(slide, notes[4])
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "A4 공유 레일 Blocking 검증", "A4 long window와 normal window를 비교해 레일 영향 확인", "05")
    add_image_fit(slide, OUT_FIG / "a4_long_vs_normal.png", 0.5, 1.12, 5.9, 3.25)
    add_image_fit(slide, OUT_FIG / "a4_cross_family_overlap.png", 6.85, 1.12, 5.7, 3.25)
    a4_rows = [[r.metric, fmt_num(float(r.value), 2) if str(r.value).replace('.', '', 1).isdigit() else r.value] for r in a4_evidence.itertuples()]
    add_table(slide, 0.75, 4.72, 5.25, 1.25, ["지표", "값"], a4_rows, 8.3)
    add_bullets(
        slide,
        6.55,
        4.65,
        5.8,
        1.35,
        [
            "A4는 사상, 레이저마킹, 로봇, 툴체인저, 틸팅이 결합된 복합 CELL이다.",
            f"현재 A4 blocking 근거 강도는 {a4_strength}로 분류했다.",
            "TRN_DEV_ID 물리 위치 매핑이 있어야 blocking 원인을 확정할 수 있다.",
        ],
        10.6,
    )
    notes[5] = "A4는 단순 설비가 아니라 여러 작업과 로봇, 툴체인저가 묶인 복합 CELL입니다. 그래서 처리시간이 길어질 가능성이 구조적으로 있습니다. 이번 분석에서는 A4 long window와 normal window를 비교하고, 해당 구간에서 다른 제품군의 레일 이벤트가 얼마나 겹치는지 봤습니다. 현재는 공유 레일 blocking 후보로 볼 근거가 있지만, 실제 물리 구간 매핑이 없기 때문에 확정 표현은 피해야 합니다."
    add_notes(slide, notes[5])
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "제품군별 병목 + 레일/AMR 자원 분석", "제품 route 차이와 공유 이송 자원을 분리해 해석", "06")
    add_image_fit(slide, OUT_FIG / "product_family_bottleneck_comparison.png", 0.5, 1.08, 5.85, 3.45)
    add_image_fit(slide, OUT_FIG / "rail_device_utilization.png", 6.85, 1.08, 5.65, 3.45)
    handle_top = ", ".join(family_rank[family_rank["product_family"] == "HANDLE"].head(3)["process"].tolist())
    mirror_top = ", ".join(family_rank[family_rank["product_family"] == "ROOM_MIRROR"].head(3)["process"].tolist())
    amr_rows = [[r.resource, int(r.bottleneck_rank), fmt_sec(r.top_related_wait_p90_sec), f"{r.alarm_ratio:.3f}"] for r in amr.itertuples()]
    add_table(slide, 0.75, 4.78, 5.25, 1.15, ["AMR", "순위", "관련 대기 p90", "alarm"], amr_rows, 8.3)
    add_bullets(
        slide,
        6.55,
        4.65,
        5.8,
        1.38,
        [
            f"HANDLE 상위 후보: {handle_top}",
            f"ROOM_MIRROR 상위 후보: {mirror_top}",
            "레일/AMR은 병목 원인 확정보다 지연 증폭 자원으로 관리해야 한다.",
        ],
        10.6,
    )
    notes[6] = "제품군을 나누면 병목 해석이 더 명확해집니다. 핸들 계열은 A4와 A6, A9_2를 포함하는 긴 흐름이고, 룸미러 계열은 A7에서 A5, A9_1로 이어지는 흐름입니다. 레일과 AMR은 모든 제품군이 공유하기 때문에 한쪽 제품군의 지연이 다른 제품군 대기로 이어질 수 있습니다. 다만 TRN_DEV_ID의 실제 위치가 없어서 레일별 물리 병목은 후보 수준으로 봅니다."
    add_notes(slide, notes[6])
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "현재 결론과 개선 단계 진입 조건", "1~5번 진단 완료 후 6번부터 rule-based scheduling 설계", "07")
    add_table(
        slide,
        0.55,
        1.15,
        6.4,
        2.25,
        ["구분", "결론"],
        [
            ["A3", "후단 반출 병목. AMR/C2/BUFFER 세분화 필요"],
            ["A4", f"공유 레일 blocking 후보. 근거 강도 {a4_strength}"],
            ["제품군", "핸들/룸미러 route 차이로 병목 분리 필요"],
            ["레일/AMR", "분석 가능하나 물리 위치/dispatch 로그 필요"],
        ],
        8.7,
    )
    add_table(
        slide,
        7.25,
        1.15,
        5.1,
        2.25,
        ["필요 정보", "이유"],
        [
            ["TRN_DEV_ID 위치", "레일 물리 blocking 확정"],
            ["AMR 배정 로그", "호출/배정 지연 분리"],
            ["C2 slot capacity", "창고 포화 판단"],
            ["OEE IDLE 정의", "정상 대기/비가동 구분"],
        ],
        8.1,
    )
    add_bullets(
        slide,
        0.85,
        4.05,
        11.2,
        1.35,
        [
            "다음 단계는 A3 우선 반출, AMR 배정, A4 WIP cap, rail scheduling 규칙을 baseline으로 설계하는 것이다.",
            "효과 검증은 기존 로그 리플레이 또는 시뮬레이션으로 대기시간/WIP/throughput 변화를 비교한다.",
            "LLM decision agent는 rule baseline 이후 state, output JSON, constraint verifier를 정의해 비교 평가한다.",
        ],
        11.4,
    )
    add_callout(slide, 0.75, 5.85, 11.7, 0.58, "최종 메시지", "현재 데이터로 1~5번 진단은 가능하며, 원인 확정에 부족한 정보는 별도 gap report로 분리했다.", BLUE)
    notes[7] = "마지막 장은 현재 결론과 다음 단계입니다. 1~5번 분석으로 병목 후보와 원인 가설은 정리됐습니다. A3는 후단 반출 병목이고, A4는 공유 레일 blocking 후보이며, 제품군별 route 차이를 분리해야 합니다. 다만 TRN_DEV_ID 위치, AMR dispatch 로그, C2 slot capacity 같은 정보가 없으면 원인을 확정할 수 없습니다. 이 정보를 확보하거나 후보 수준으로 둔 상태에서 6번 rule-based scheduling baseline 설계로 넘어가면 됩니다."
    add_notes(slide, notes[7])
    add_footer(slide, 7)

    prs.save(PPT_OUT)
    return notes


def write_script(notes: dict[int, str]) -> None:
    titles = [
        "분석 목적과 현재 완료 범위",
        "데이터 검증과 정상 Route 기준",
        "전체 병목 스크리닝 결과",
        "A3 병목 세분화",
        "A4 공유 레일 Blocking 검증",
        "제품군별 병목 + 레일/AMR 자원 분석",
        "현재 결론과 개선 단계 진입 조건",
    ]
    lines = ["# P-ZONE 1~5번 분석 PPT 발표 대본", ""]
    for idx, title in enumerate(titles, start=1):
        lines.append(f"## {idx}. {title}")
        lines.append("")
        lines.append(notes[idx])
        lines.append("")
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")


def load_all_data() -> dict[str, pd.DataFrame]:
    return {
        "all_routes": read_csv("all_product_routes_before_filter.csv"),
        "routes": read_csv("product_routes.csv"),
        "events": read_csv("equipment_processing_events.csv"),
        "wait_events": read_csv("transition_waiting_events.csv"),
        "pairing_anomalies": read_csv("pairing_anomalies.csv"),
        "normal_route_summary": read_csv("normal_route_summary.csv"),
        "excluded_routes": read_csv("excluded_route_summary.csv"),
        "wip_summary": read_csv("wip_summary.csv"),
        "bottleneck": read_csv("bottleneck_ranking.csv", PKG_DATA),
        "wait_summary": read_csv("transition_waiting_time.csv", PKG_DATA),
        "occupancy": read_csv("occupancy_summary.csv", PKG_DATA),
        "oee": read_csv("oee_summary.csv", PKG_DATA),
        "a4_windows": read_csv("a4_blocking_windows.csv", PKG_DATA),
        "a4_downstream": read_csv("a4_downstream_lag.csv", PKG_DATA),
    }


def run() -> None:
    ensure_dirs()
    data = load_all_data()
    sns.set_theme(style="whitegrid", font=FONT)
    results: dict[str, dict[str, pd.DataFrame]] = {
        "base": {"bottleneck": data["bottleneck"]},
        "validation": build_data_validation(data),
        "a3": build_a3_analysis(data),
        "a4": build_a4_analysis(data),
        "family": build_product_family_analysis(data),
        "transport": build_transport_analysis(data),
    }
    write_required_gap_report(results["transport"]["limitations"])
    make_markdown_summary(results)
    notes = build_ppt(results)
    write_script(notes)
    print(f"saved analysis package: {OUT}")
    print(f"saved ppt: {PPT_OUT}")
    print(f"saved script: {SCRIPT_OUT}")


if __name__ == "__main__":
    run()
