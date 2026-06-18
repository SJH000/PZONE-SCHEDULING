from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("outputs/2026_06_14")
DATA = BASE / "data"
FIG = BASE / "figures"
OUT = BASE / "PZONE_bottleneck_summary_7slides.pptx"

FONT = "Malgun Gothic"
NAVY = RGBColor(24, 39, 61)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(13, 148, 136)
ORANGE = RGBColor(234, 88, 12)
GRAY = RGBColor(75, 85, 99)
LIGHT_GRAY = RGBColor(243, 244, 246)
MID_GRAY = RGBColor(209, 213, 219)
WHITE = RGBColor(255, 255, 255)


def set_text_style(shape, font_size=14, color=NAVY, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(slide, x, y, w, h, text, font_size=14, color=NAVY, bold=False, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    set_text_style(shape, font_size, color, bold)
    return shape


def add_title(slide, title, subtitle=None, section=None):
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, 0.45, 0.35, 9.8, 0.45, title, 24, NAVY, True)
    if subtitle:
        add_textbox(slide, 0.48, 0.82, 10.6, 0.32, subtitle, 10.5, GRAY)
    if section:
        add_textbox(slide, 11.6, 0.42, 1.25, 0.3, section, 9.5, GRAY, False, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, font_size=13, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_callout(slide, x, y, w, h, title, body, color=BLUE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(239, 246, 255)
    rect.line.color.rgb = color
    rect.line.width = Pt(1)
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.24, title, 10.5, color, True)
    add_textbox(slide, x + 0.18, y + 0.43, w - 0.36, h - 0.55, body, 11.5, NAVY)


def add_image_fit(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio >= box_ratio:
        final_w = w
        final_h = w / ratio
    else:
        final_h = h
        final_w = h * ratio
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(final_w), height=Inches(final_h))


def add_table(slide, x, y, w, h, headers, rows, font_size=9.5, header_fill=LIGHT_GRAY):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = NAVY
    for cell in table.rows[0].cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
    return table_shape


def add_footer(slide, page):
    add_textbox(slide, 0.45, 7.12, 5.2, 0.2, "P-ZONE 병목 분석 | 정상 제품/정상 route 기준", 8.5, GRAY)
    add_textbox(slide, 12.3, 7.12, 0.55, 0.2, f"{page}/7", 8.5, GRAY, False, PP_ALIGN.RIGHT)


def sec(value):
    if pd.isna(value):
        return "-"
    return f"{float(value):,.1f}s"


def pct(value):
    if pd.isna(value):
        return "-"
    return f"{float(value) * 100:.1f}%"


def load_data():
    return {
        "routes": pd.read_csv(DATA / "normal_route_summary.csv"),
        "bottlenecks": pd.read_csv(DATA / "bottleneck_ranking.csv"),
        "waiting": pd.read_csv(DATA / "transition_waiting_time.csv"),
        "processing": pd.read_csv(DATA / "equipment_processing_time_summary.csv"),
        "sched": pd.read_csv(DATA / "schedulability_classification.csv"),
        "wip": pd.read_csv(DATA / "wip_summary.csv"),
        "a4": pd.read_csv(DATA / "a4_blocking_windows.csv"),
        "downstream": pd.read_csv(DATA / "a4_downstream_lag.csv"),
    }


def make_slide_1(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "분석 목적과 데이터 범위", "전체 공정을 먼저 같은 기준으로 스크리닝한 뒤, 근거가 강한 병목만 심층 분석", "01")

    normal_count = 294
    total_count = 324
    excluded_count = 30
    table_rows = [
        ["분석 기준", "정상 제품 + 정상 route/prefix serial"],
        ["제품군", "PRD1000, PRD2000, PRD3001, PRD3002"],
        ["정상/전체", f"{normal_count}/{total_count} serial"],
        ["제외", f"{excluded_count} serial: 비정상 route, 음수/누락/극단 이상치"],
        ["핵심 테이블", "prc_hist, prc_trns, prc_oee, strg_buf, strg_fns, std_eqpmn"],
    ]
    add_table(slide, 0.6, 1.45, 5.9, 2.2, ["항목", "내용"], table_rows, 10.5)
    add_callout(
        slide,
        7.0,
        1.45,
        5.5,
        1.45,
        "핵심 질문",
        "어느 공정이 실제 병목이며, 그 원인이 스케줄링으로 완화 가능한지 또는 물리적 제약인지 구분한다.",
        BLUE,
    )
    add_bullets(
        slide,
        7.05,
        3.18,
        5.2,
        2.0,
        [
            "처리시간, 전이 대기시간, WIP, 이송/레일 점유, OEE 상태를 통합했다.",
            "매뉴얼 기반 공정 의미를 함께 사용해 숫자의 원인을 해석했다.",
            "A3/A4/A9/A7/A5/AMR은 전체 스크리닝에서 근거가 강해 심층 분석 대상으로 선정했다.",
        ],
        13,
    )
    add_callout(slide, 0.6, 5.18, 11.9, 0.9, "현재 결론 한 줄", "A3는 후단 배출 운영 병목, A4는 복합 CELL 처리시간과 공유 레일 blocking 후보, A9/A7/A5/AMR은 연결 병목 후보로 판단된다.", ORANGE)
    add_footer(slide, 1)


def make_slide_2(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "제품별 정상 Route와 Clean 기준", "A5 -> A3 같은 예외 전이를 제거하고 정상 흐름만 분석", "02")

    routes = data["routes"].copy()
    rows = []
    for _, row in routes.iterrows():
        route = row["normal_route"].replace("C1_RAW_STORAGE", "C1").replace("C2_FINISHED_STORAGE", "C2")
        if len(route) > 82:
            route = route[:79] + "..."
        rows.append([row["PRD_CD"], route, int(row["support_count_exact"])])
    add_table(slide, 0.55, 1.25, 12.25, 2.35, ["제품", "정상 route", "정확 일치 수"], rows, 8.6)

    add_bullets(
        slide,
        0.72,
        3.92,
        5.75,
        1.7,
        [
            "핸들 계열 PRD1000/PRD2000은 A4와 A6를 포함하는 긴 route를 가진다.",
            "룸미러 계열 PRD3001/PRD3002는 A7 -> A5 -> A9_1 -> A3 흐름이 정상이다.",
            "따라서 A7 -> A5는 예외가 아니라 제품군별 정상 route 차이다.",
        ],
        12.5,
    )
    add_bullets(
        slide,
        6.95,
        3.92,
        5.65,
        1.7,
        [
            "Clean 기준은 이상치 제외 후 비교 가능한 지표를 뜻한다.",
            "처리/대기시간이 음수이거나 누락된 이벤트는 anomaly로 분리했다.",
            "설비별 Q3 + 1.5*IQR 초과 또는 600초 초과는 raw와 clean을 분리해 해석했다.",
        ],
        12.5,
    )
    add_callout(slide, 0.6, 5.9, 11.9, 0.55, "해석 포인트", "정상 route 기준을 먼저 세우지 않으면 제품군별 정상 차이와 비정상 로그가 섞여 병목 위치가 왜곡된다.", TEAL)
    add_footer(slide, 2)


def make_slide_3(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "전체 공정 병목 스크리닝", "모든 공정을 같은 지표로 점수화한 뒤 우선순위를 정렬", "03")
    add_image_fit(slide, FIG / "01_bottleneck_ranking.png", 0.55, 1.17, 7.5, 4.9)

    b = data["bottlenecks"].head(7)
    rows = [[int(r["rank"]), r["process"], f"{r['bottleneck_score']:.2f}", sec(r["processing_p90_sec"]), sec(r["waiting_p90_sec"])] for _, r in b.iterrows()]
    add_table(slide, 8.25, 1.3, 4.65, 2.15, ["순위", "공정", "점수", "처리 p90", "대기 p90"], rows, 7.8)
    add_bullets(
        slide,
        8.35,
        3.75,
        4.35,
        1.95,
        [
            "A3가 압도적으로 높다: 처리보다 후단 대기/WIP 영향이 크다.",
            "A4는 처리 p90과 점유시간이 높아 구조적 부하 후보다.",
            "AMR_LD90, A7, A9_1, A5, A9_2는 연결 병목 후보로 남긴다.",
        ],
        11.5,
    )
    add_footer(slide, 3)


def make_slide_4(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "A3 후단 배출 병목", "A3 자체 처리보다 A3 이후 AMR/C2 반출 대기가 병목 신호", "04")
    add_image_fit(slide, FIG / "03_waiting_p90_top.png", 0.55, 1.12, 7.0, 4.75)

    wait = data["waiting"]
    a3_c2 = wait[wait["transition"] == "A3->C2_FINISHED_STORAGE"].iloc[0]
    a3_amr = wait[wait["transition"] == "A3->AMR_LD250"].iloc[0]
    wip = data["wip"]
    buffer_wip = wip[wip["area"] == "BUFFER"]["p90_wip"].iloc[0]
    rows = [
        ["A3 처리 p90", sec(data["bottlenecks"].query("process == 'A3'")["processing_p90_sec"].iloc[0])],
        ["A3 -> AMR_LD250 대기 p90(clean)", sec(a3_amr["p90_wait_sec_clean"])],
        ["A3 -> C2 대기 p90(clean)", sec(a3_c2["p90_wait_sec_clean"])],
        ["BUFFER WIP p90", f"{buffer_wip:.1f}"],
    ]
    add_table(slide, 7.9, 1.25, 4.75, 1.55, ["지표", "값"], rows, 10)
    add_bullets(
        slide,
        7.95,
        3.08,
        4.5,
        2.0,
        [
            "A3는 완제품/NG 분류, 2단 적재 대기, 스토퍼 기반 1개씩 배출 구조다.",
            "대기 p90이 처리 p90보다 훨씬 커서 설비 가공보다 후단 반출 운영 문제가 크다.",
            "우선 액션은 PRIORITIZE_DISCHARGE와 DISPATCH_AMR이다.",
        ],
        11.5,
    )
    add_callout(slide, 7.9, 5.45, 4.75, 0.75, "판단", "스케줄링 개선 가능성이 가장 큰 병목이다. 단, C2 적재/AMR 가용성 확인이 필요하다.", ORANGE)
    add_footer(slide, 4)


def make_slide_5(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "A4 복합 CELL과 공유 레일 Blocking 후보", "A4 long window가 공유 레일 흐름에 영향을 주는지 검증", "05")
    add_image_fit(slide, FIG / "02_processing_p90_top.png", 0.45, 1.08, 5.95, 3.0)
    add_image_fit(slide, FIG / "05_a4_processing_vs_rail_occupancy.png", 6.75, 1.05, 5.95, 3.0)

    a4_proc = data["processing"].query("logical_process == 'A4'").iloc[0]
    a4_windows = data["a4"]
    rows = [
        ["A4 처리 p90(clean)", sec(a4_proc["p90_sec_clean"])],
        ["A4 이상치 제외 후 count", int(a4_proc["count_clean"])],
        ["A4 long window 수", len(a4_windows)],
        ["long 기준", sec(a4_windows["a4_long_threshold_sec"].iloc[0])],
    ]
    add_table(slide, 0.75, 4.42, 4.35, 1.45, ["지표", "값"], rows, 9.4)
    add_bullets(
        slide,
        5.45,
        4.32,
        6.75,
        1.65,
        [
            "매뉴얼상 A4는 사상 제거, 레이저 마킹, 산업용/협동 로봇, 툴체인저, 틸팅, 3개 파렛트 버퍼가 결합된 복합 공정이다.",
            "A4 처리시간이 길어지면 공유 레일 통과가 지연되고, 다른 제품군의 이송 이벤트와 겹칠 가능성이 생긴다.",
            "우선 액션은 CONTROL_WIP, SCHEDULE_RAIL, HOLD_ENTRY 조합이다.",
        ],
        10.8,
    )
    add_footer(slide, 5)


def make_slide_6(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "연결 병목: A9 / A7 / A5 / AMR", "제품군 route 차이와 후단 반출 자원이 병목을 연결한다", "06")
    add_image_fit(slide, FIG / "04_rail_occupancy_timeseries.png", 0.55, 1.08, 6.35, 3.1)
    add_image_fit(slide, FIG / "06_a4_blocking_overlap_by_product_family.png", 7.05, 1.08, 5.55, 3.1)

    b = data["bottlenecks"]
    targets = ["AMR_LD90", "A7", "A9_1", "A5", "A9_2", "AMR_LD250"]
    rows = []
    for process in targets:
        row = b[b["process"] == process]
        if row.empty:
            continue
        r = row.iloc[0]
        rows.append([process, int(r["rank"]), sec(r["processing_p90_sec"]), sec(r["waiting_p90_sec"])])
    add_table(slide, 0.65, 4.48, 5.7, 1.55, ["공정", "순위", "처리 p90", "대기 p90"], rows, 8.7)
    add_bullets(
        slide,
        6.8,
        4.45,
        5.75,
        1.55,
        [
            "A9_1/A9_2는 검사·조립 후단이며 A3 반출 흐름과 연결된다.",
            "A7/A5는 제품군별 route가 달라 정상 전이 해석이 중요하다.",
            "AMR_LD90/LD250은 레일/반출 자원으로 공정 간 지연을 증폭시킬 수 있다.",
        ],
        10.9,
    )
    add_footer(slide, 6)


def make_slide_7(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "스케줄링 가능성 및 후속 과제", "즉시 조정 가능한 운영 병목과 물리 제약 후보를 분리", "07")

    sched = data["sched"].head(7)
    rows = []
    for _, r in sched.iterrows():
        rows.append([int(r["rank"]), r["process"], r["schedulability"], r["recommended_action"]])
    add_table(slide, 0.55, 1.15, 6.9, 2.35, ["순위", "공정", "분류", "추천 액션"], rows, 8.5)
    add_image_fit(slide, FIG / "07_a4_downstream_lag_effect.png", 7.75, 1.12, 4.85, 3.0)

    add_bullets(
        slide,
        0.75,
        3.88,
        5.9,
        1.85,
        [
            "A3: 반출 우선순위, AMR 배정, C2 적재 가능 여부를 분리 검증한다.",
            "A4: long/normal window 대조군으로 공유 레일 blocking을 강화 검증한다.",
            "제품군별 병목 순위를 핸들 계열과 룸미러 계열로 분리한다.",
        ],
        11.2,
    )
    add_bullets(
        slide,
        7.15,
        4.35,
        5.25,
        1.65,
        [
            "Rule baseline: PRIORITIZE_DISCHARGE, DISPATCH_AMR, CONTROL_WIP, SCHEDULE_RAIL, HOLD_ENTRY.",
            "다음 단계: state 정의, constraint verifier, LLM-assisted Decision Agent 비교 평가.",
        ],
        11.2,
    )
    add_callout(slide, 0.7, 6.05, 11.75, 0.55, "최종 메시지", "현재 데이터 기준 1순위 개선 대상은 A3 후단 반출, 2순위 검증 대상은 A4 공유 레일 blocking이다.", BLUE)
    add_footer(slide, 7)


def build():
    data = load_data()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_slide_1(prs, data)
    make_slide_2(prs, data)
    make_slide_3(prs, data)
    make_slide_4(prs, data)
    make_slide_5(prs, data)
    make_slide_6(prs, data)
    make_slide_7(prs, data)

    prs.save(OUT)
    print(f"saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
